from __future__ import annotations

import io
import json
import os
import re
from contextlib import suppress
from datetime import date, datetime, timedelta
from decimal import Decimal
from typing import Any, Dict, List, Optional, Sequence, Tuple

from backend.database.dbconnect import create_metadata_connection
from backend.modules.common.db_table_utils import _detect_db_type, get_postgresql_table_name
from backend.modules.reports.report_service import ReportMetadataService, ReportServiceError


class DashboardCreatorError(Exception):
    def __init__(
        self,
        message: str,
        status_code: int = 400,
        code: str = "DASHBOARD_CREATOR_ERROR",
        details: Optional[Dict[str, Any]] = None,
    ):
        super().__init__(message)
        self.message = message
        self.status_code = status_code
        self.code = code
        self.details = details or {}


class _ParamBuilder:
    def __init__(self, db_type: str):
        self.db_type = (db_type or "ORACLE").upper()
        self._counter = 0
        if self.db_type == "ORACLE":
            self._params: Dict[str, Any] | List[Any] = {}
        else:
            self._params = []

    def add(self, value: Any, hint: str = "p") -> str:
        if self.db_type == "ORACLE":
            key = f"{hint}{self._counter}"
            self._counter += 1
            self._params[key] = value
            return f":{key}"
        self._params.append(value)
        return "%s"

    @property
    def params(self) -> Dict[str, Any] | Sequence[Any] | None:
        if self.db_type == "ORACLE":
            return self._params if self._params else None
        return tuple(self._params) if self._params else None


class DashboardCreatorService:
    def __init__(self):
        self.schema = (os.getenv("DMS_SCHEMA") or "").strip()
        self.report_service = ReportMetadataService()
        self.max_widgets = 25
        self.max_sql_length = 120000
        self.max_name_length = 200
        self.max_description_length = 1000
        self.allowed_widget_types = {"TABLE", "BAR", "LINE", "PIE", "KPI", "AREA"}

    def _table_ref(self, cursor, db_type: str, base_table_name: str) -> str:
        if db_type == "POSTGRESQL":
            schema_lower = self.schema.lower() if self.schema else "public"
            actual_name = get_postgresql_table_name(cursor, schema_lower, base_table_name)
            table_ref = (
                f'"{actual_name}"' if actual_name != actual_name.lower() else actual_name
            )
            return f"{schema_lower}.{table_ref}"

        schema_prefix = f"{self.schema}." if self.schema else ""
        return f"{schema_prefix}{base_table_name}"

    def _open_connection(self):
        conn = create_metadata_connection()
        cursor = conn.cursor()
        db_type = _detect_db_type(conn)
        tables = {
            "dash_def": self._table_ref(cursor, db_type, "DMS_DASH_DEF"),
            "dash_widget": self._table_ref(cursor, db_type, "DMS_DASH_WIDGET"),
            "dash_filter": self._table_ref(cursor, db_type, "DMS_DASH_FILTER"),
            "dash_share": self._table_ref(cursor, db_type, "DMS_DASH_SHARE"),
            "dash_export_log": self._table_ref(cursor, db_type, "DMS_DASH_EXPORT_LOG"),
        }
        return conn, cursor, db_type, tables

    def _close_connection(self, conn, cursor):
        with suppress(Exception):
            if cursor:
                cursor.close()
        with suppress(Exception):
            if conn:
                conn.close()

    def _execute(self, cursor, query: str, params: Optional[Sequence[Any] | Dict[str, Any]] = None):
        if params is None:
            cursor.execute(query)
        else:
            cursor.execute(query, params)

    def _fetch_all_dict(self, cursor) -> List[Dict[str, Any]]:
        rows = cursor.fetchall()
        if not rows:
            return []
        columns = [desc[0].lower() for desc in cursor.description]
        return [
            {columns[index]: row[index] for index in range(len(columns))}
            for row in rows
        ]

    def _read_lob(self, value: Any):
        if value is None:
            return None
        if hasattr(value, "read"):
            try:
                return value.read()
            except Exception:
                return str(value)
        return value

    def _to_int(self, value: Any) -> Optional[int]:
        if value is None:
            return None
        try:
            return int(value)
        except Exception:
            return None

    def _to_flag(self, value: bool) -> str:
        return "Y" if bool(value) else "N"

    def _from_flag(self, value: Any, default: bool = True) -> bool:
        if value is None:
            return default
        return str(value).upper() in {"Y", "1", "TRUE", "T"}

    def _to_iso(self, value: Any) -> Optional[str]:
        if value is None:
            return None
        if hasattr(value, "isoformat"):
            return value.isoformat()
        return str(value)

    def _json_safe(self, value: Any):
        if isinstance(value, Decimal):
            return float(value)
        if isinstance(value, timedelta):
            return value.total_seconds()
        if isinstance(value, (datetime, date)):
            return value.isoformat()
        return value

    def _next_id(self, cursor, sequence_name: str) -> int:
        cursor.execute(f"SELECT {sequence_name}.NEXTVAL FROM DUAL")
        row = cursor.fetchone()
        if not row:
            raise DashboardCreatorError("Failed to generate id", code="ID_GENERATION_FAILED")
        return int(row[0])

    def _validate_read_only_sql(self, sql_text: str):
        sql_clean = (sql_text or "").strip()
        if not sql_clean:
            raise DashboardCreatorError("SQL text is required", code="SQL_TEXT_REQUIRED")
        if len(sql_clean) > self.max_sql_length:
            raise DashboardCreatorError(
                "SQL text exceeds allowed size",
                code="SQL_TEXT_TOO_LARGE",
                details={"maxLength": self.max_sql_length},
            )

        normalized = sql_clean.lower()
        if not (normalized.startswith("select") or normalized.startswith("with")):
            raise DashboardCreatorError(
                "Only SELECT/CTE statements are allowed",
                code="SQL_NOT_SELECT",
            )

        blocked_keywords = [
            " insert ",
            " update ",
            " delete ",
            " drop ",
            " alter ",
            " truncate ",
            " create ",
            " grant ",
            " revoke ",
            " merge ",
            " execute ",
            " commit",
            " rollback",
        ]
        padded = f" {normalized} "
        for keyword in blocked_keywords:
            if keyword in padded:
                raise DashboardCreatorError(
                    "Only read-only SQL is allowed",
                    code="SQL_NOT_READ_ONLY",
                    details={"keyword": keyword.strip()},
                )

    def _normalize_dashboard_name(self, dashboard_name: Any) -> str:
        normalized = str(dashboard_name or "").strip()
        if not normalized:
            raise DashboardCreatorError("dashboardName is required", code="DASHBOARD_NAME_REQUIRED")
        if len(normalized) > self.max_name_length:
            raise DashboardCreatorError(
                "dashboardName is too long",
                code="DASHBOARD_NAME_TOO_LONG",
                details={"maxLength": self.max_name_length},
            )
        return normalized

    def _normalize_description(self, description: Any) -> Optional[str]:
        if description is None:
            return None
        normalized = str(description)
        if len(normalized) > self.max_description_length:
            raise DashboardCreatorError(
                "description is too long",
                code="DASHBOARD_DESCRIPTION_TOO_LONG",
                details={"maxLength": self.max_description_length},
            )
        return normalized

    def _normalize_widgets(self, widgets_payload: Any) -> List[Dict[str, Any]]:
        widgets = widgets_payload or []
        if not isinstance(widgets, list):
            raise DashboardCreatorError("widgets must be a list", code="INVALID_WIDGETS_PAYLOAD")
        if len(widgets) > self.max_widgets:
            raise DashboardCreatorError(
                "Too many widgets in dashboard",
                code="WIDGET_LIMIT_EXCEEDED",
                details={"maxWidgets": self.max_widgets},
            )

        normalized_widgets: List[Dict[str, Any]] = []
        for index, widget in enumerate(widgets, start=1):
            if not isinstance(widget, dict):
                raise DashboardCreatorError(
                    "Each widget must be an object",
                    code="INVALID_WIDGET",
                    details={"index": index},
                )

            widget_name = str(widget.get("widgetName") or f"Widget {index}").strip()
            if len(widget_name) > self.max_name_length:
                raise DashboardCreatorError(
                    "widgetName is too long",
                    code="WIDGET_NAME_TOO_LONG",
                    details={"index": index, "maxLength": self.max_name_length},
                )

            widget_type = str(widget.get("widgetType") or "TABLE").upper()
            if widget_type not in self.allowed_widget_types:
                raise DashboardCreatorError(
                    "Unsupported widgetType",
                    code="INVALID_WIDGET_TYPE",
                    details={"index": index, "widgetType": widget_type},
                )

            source_mode = str(widget.get("sourceMode") or "SQL").upper()
            if source_mode not in {"SQL", "TABLE", "REPORT_REF"}:
                raise DashboardCreatorError(
                    "Unsupported sourceMode",
                    code="INVALID_SOURCE_MODE",
                    details={"index": index, "sourceMode": source_mode},
                )

            adhoc_sql = widget.get("adhocSql")
            sql_source_id = self._to_int(widget.get("sqlSourceId"))

            if source_mode == "SQL":
                if not adhoc_sql and not sql_source_id:
                    raise DashboardCreatorError(
                        "Widget SQL source is required",
                        code="WIDGET_SQL_REQUIRED",
                        details={"index": index},
                    )
                if adhoc_sql:
                    self._validate_read_only_sql(str(adhoc_sql))

            db_connection_id = self._to_int(widget.get("dbConnectionId"))

            normalized_widgets.append(
                {
                    "widgetName": widget_name,
                    "widgetType": widget_type,
                    "sourceMode": source_mode,
                    "sqlSourceId": sql_source_id,
                    "adhocSql": adhoc_sql,
                    "dbConnectionId": db_connection_id,
                    "configJson": widget.get("configJson"),
                    "layoutJson": widget.get("layoutJson"),
                    "isActive": bool(widget.get("isActive", True)),
                }
            )

        return normalized_widgets

    def _serialize_dashboard(self, row: Dict[str, Any]) -> Dict[str, Any]:
        return {
            "dashboardId": self._to_int(row.get("dashid")),
            "dashboardName": row.get("dashnm"),
            "description": self._read_lob(row.get("dscrptn")),
            "ownerUserId": self._to_int(row.get("owner_usrid")),
            "isActive": self._from_flag(row.get("is_actv"), default=True),
            "curFlg": row.get("curflg"),
            "createdBy": row.get("crtdby"),
            "createdAt": self._to_iso(row.get("crtddt")),
            "updatedBy": row.get("updtdby"),
            "updatedAt": self._to_iso(row.get("updtdt")),
            "widgetCount": self._to_int(row.get("widget_count")) or 0,
        }

    def _serialize_widget(self, row: Dict[str, Any]) -> Dict[str, Any]:
        return {
            "widgetId": self._to_int(row.get("widgtid")),
            "dashboardId": self._to_int(row.get("dashid")),
            "widgetName": row.get("widgtnm"),
            "widgetType": row.get("widgttyp"),
            "sourceMode": row.get("srcmode"),
            "sqlSourceId": self._to_int(row.get("sqlsrcid")),
            "adhocSql": self._read_lob(row.get("adhcsql")),
            "dbConnectionId": self._to_int(row.get("dbcnid")),
            "configJson": self._read_lob(row.get("cfg_json")),
            "layoutJson": self._read_lob(row.get("layout_json")),
            "orderNo": self._to_int(row.get("order_no")) or 1,
            "isActive": self._from_flag(row.get("is_actv"), default=True),
            "curFlg": row.get("curflg"),
        }

    def list_dashboards(
        self,
        search: Optional[str] = None,
        include_inactive: bool = False,
        owner_user_id: Optional[int] = None,
    ) -> List[Dict[str, Any]]:
        conn, cursor, db_type, tables = self._open_connection()
        try:
            builder = _ParamBuilder(db_type)
            query = f"""
                SELECT
                    d.DASHID AS dashid,
                    d.DASHNM AS dashnm,
                    d.DSCRPTN AS dscrptn,
                    d.OWNER_USRID AS owner_usrid,
                    d.IS_ACTV AS is_actv,
                    d.CURFLG AS curflg,
                    d.CRTDBY AS crtdby,
                    d.CRTDDT AS crtddt,
                    d.UPDTDBY AS updtdby,
                    d.UPDTDT AS updtdt,
                    (
                        SELECT COUNT(1)
                        FROM {tables['dash_widget']} w
                        WHERE w.DASHID = d.DASHID
                          AND w.CURFLG = 'Y'
                    ) AS widget_count
                FROM {tables['dash_def']} d
                WHERE d.CURFLG = 'Y'
            """

            if not include_inactive:
                query += " AND d.IS_ACTV = 'Y'"

            if owner_user_id is not None:
                query += f" AND d.OWNER_USRID = {builder.add(owner_user_id, 'owner')}"

            if search:
                pattern = f"%{search.upper()}%"
                p1 = builder.add(pattern, "search")
                p2 = builder.add(pattern, "search")
                query += f" AND (UPPER(d.DASHNM) LIKE {p1} OR UPPER(d.DSCRPTN) LIKE {p2})"

            query += " ORDER BY d.DASHNM"
            self._execute(cursor, query, builder.params)
            rows = self._fetch_all_dict(cursor)
            return [self._serialize_dashboard(row) for row in rows]
        finally:
            self._close_connection(conn, cursor)

    def get_dashboard(self, dashboard_id: int) -> Dict[str, Any]:
        conn, cursor, db_type, tables = self._open_connection()
        try:
            builder = _ParamBuilder(db_type)
            id_placeholder = builder.add(dashboard_id, "dash")
            query = f"""
                SELECT
                    d.DASHID AS dashid,
                    d.DASHNM AS dashnm,
                    d.DSCRPTN AS dscrptn,
                    d.OWNER_USRID AS owner_usrid,
                    d.IS_ACTV AS is_actv,
                    d.CURFLG AS curflg,
                    d.CRTDBY AS crtdby,
                    d.CRTDDT AS crtddt,
                    d.UPDTDBY AS updtdby,
                    d.UPDTDT AS updtdt
                FROM {tables['dash_def']} d
                WHERE d.DASHID = {id_placeholder}
                  AND d.CURFLG = 'Y'
            """
            self._execute(cursor, query, builder.params)
            rows = self._fetch_all_dict(cursor)
            if not rows:
                raise DashboardCreatorError(
                    "Dashboard not found",
                    status_code=404,
                    code="DASHBOARD_NOT_FOUND",
                )
            dashboard = self._serialize_dashboard(rows[0])

            widget_builder = _ParamBuilder(db_type)
            widget_id_placeholder = widget_builder.add(dashboard_id, "dash")
            widget_query = f"""
                SELECT
                    w.WIDGTID AS widgtid,
                    w.DASHID AS dashid,
                    w.WIDGTNM AS widgtnm,
                    w.WIDGTTYP AS widgttyp,
                    w.SRCMODE AS srcmode,
                    w.SQLSRCID AS sqlsrcid,
                    w.ADHCSQL AS adhcsql,
                    w.DBCNID AS dbcnid,
                    w.CFG_JSON AS cfg_json,
                    w.LAYOUT_JSON AS layout_json,
                    w.ORDER_NO AS order_no,
                    w.IS_ACTV AS is_actv,
                    w.CURFLG AS curflg
                FROM {tables['dash_widget']} w
                WHERE w.DASHID = {widget_id_placeholder}
                  AND w.CURFLG = 'Y'
                ORDER BY w.ORDER_NO, w.WIDGTID
            """
            self._execute(cursor, widget_query, widget_builder.params)
            widget_rows = self._fetch_all_dict(cursor)
            dashboard["widgets"] = [self._serialize_widget(row) for row in widget_rows]
            dashboard["widgetCount"] = len(dashboard["widgets"])
            return dashboard
        finally:
            self._close_connection(conn, cursor)

    def _insert_widget(
        self,
        cursor,
        db_type: str,
        tables: Dict[str, str],
        dashboard_id: int,
        widget: Dict[str, Any],
        username: str,
        order_no: int,
    ):
        if db_type == "ORACLE":
            widget_id = self._next_id(cursor, "DMS_DASH_WIDGET_SEQ")
            query = f"""
                INSERT INTO {tables['dash_widget']} (
                    WIDGTID, DASHID, WIDGTNM, WIDGTTYP, SRCMODE, SQLSRCID,
                    ADHCSQL, DBCNID, CFG_JSON, LAYOUT_JSON, ORDER_NO,
                    IS_ACTV, CURFLG, CRTDBY, CRTDDT, UPDTDBY, UPDTDT
                ) VALUES (
                    :widgtid, :dashid, :widgtnm, :widgttyp, :srcmode, :sqlsrcid,
                    :adhcsql, :dbcnid, :cfg_json, :layout_json, :order_no,
                    :is_actv, 'Y', :crtdby, SYSTIMESTAMP, :updtdby, SYSTIMESTAMP
                )
            """
            params = {
                "widgtid": widget_id,
                "dashid": dashboard_id,
                "widgtnm": widget.get("widgetName") or f"Widget {order_no}",
                "widgttyp": widget.get("widgetType") or "TABLE",
                "srcmode": widget.get("sourceMode") or "SQL",
                "sqlsrcid": self._to_int(widget.get("sqlSourceId")),
                "adhcsql": widget.get("adhocSql"),
                "dbcnid": self._to_int(widget.get("dbConnectionId")),
                "cfg_json": widget.get("configJson"),
                "layout_json": widget.get("layoutJson"),
                "order_no": order_no,
                "is_actv": self._to_flag(widget.get("isActive", True)),
                "crtdby": username,
                "updtdby": username,
            }
            self._execute(cursor, query, params)
            return

        query = f"""
            INSERT INTO {tables['dash_widget']} (
                DASHID, WIDGTNM, WIDGTTYP, SRCMODE, SQLSRCID,
                ADHCSQL, DBCNID, CFG_JSON, LAYOUT_JSON, ORDER_NO,
                IS_ACTV, CURFLG, CRTDBY, CRTDDT, UPDTDBY, UPDTDT
            ) VALUES (
                %s, %s, %s, %s, %s,
                %s, %s, %s, %s, %s,
                %s, 'Y', %s, CURRENT_TIMESTAMP, %s, CURRENT_TIMESTAMP
            )
        """
        params = (
            dashboard_id,
            widget.get("widgetName") or f"Widget {order_no}",
            widget.get("widgetType") or "TABLE",
            widget.get("sourceMode") or "SQL",
            self._to_int(widget.get("sqlSourceId")),
            widget.get("adhocSql"),
            self._to_int(widget.get("dbConnectionId")),
            widget.get("configJson"),
            widget.get("layoutJson"),
            order_no,
            self._to_flag(widget.get("isActive", True)),
            username,
            username,
        )
        self._execute(cursor, query, params)

    def create_dashboard(
        self,
        payload: Dict[str, Any],
        username: str = "system",
        owner_user_id: Optional[int] = None,
    ) -> Dict[str, Any]:
        dash_name = self._normalize_dashboard_name(payload.get("dashboardName"))
        description = self._normalize_description(payload.get("description"))
        payload_owner_user_id = self._to_int(payload.get("ownerUserId"))
        resolved_owner_user_id = owner_user_id if owner_user_id is not None else payload_owner_user_id
        is_active = self._to_flag(payload.get("isActive", True))
        widgets = self._normalize_widgets(payload.get("widgets") or [])

        conn, cursor, db_type, tables = self._open_connection()
        try:
            if db_type == "ORACLE":
                dashboard_id = self._next_id(cursor, "DMS_DASH_DEF_SEQ")
                insert_query = f"""
                    INSERT INTO {tables['dash_def']} (
                        DASHID, DASHNM, DSCRPTN, OWNER_USRID, IS_ACTV, CURFLG,
                        CRTDBY, CRTDDT, UPDTDBY, UPDTDT
                    ) VALUES (
                        :dashid, :dashnm, :dscrptn, :owner_usrid, :is_actv, 'Y',
                        :crtdby, SYSTIMESTAMP, :updtdby, SYSTIMESTAMP
                    )
                """
                params = {
                    "dashid": dashboard_id,
                    "dashnm": dash_name,
                    "dscrptn": description,
                    "owner_usrid": resolved_owner_user_id,
                    "is_actv": is_active,
                    "crtdby": username,
                    "updtdby": username,
                }
                self._execute(cursor, insert_query, params)
            else:
                insert_query = f"""
                    INSERT INTO {tables['dash_def']} (
                        DASHNM, DSCRPTN, OWNER_USRID, IS_ACTV, CURFLG,
                        CRTDBY, CRTDDT, UPDTDBY, UPDTDT
                    ) VALUES (
                        %s, %s, %s, %s, 'Y',
                        %s, CURRENT_TIMESTAMP, %s, CURRENT_TIMESTAMP
                    )
                    RETURNING DASHID
                """
                params = (dash_name, description, resolved_owner_user_id, is_active, username, username)
                self._execute(cursor, insert_query, params)
                row = cursor.fetchone()
                dashboard_id = int(row[0])

            for index, widget in enumerate(widgets, start=1):
                self._insert_widget(cursor, db_type, tables, dashboard_id, widget, username, index)

            with suppress(Exception):
                conn.commit()
            return self.get_dashboard(dashboard_id)
        except DashboardCreatorError:
            with suppress(Exception):
                conn.rollback()
            raise
        except Exception as exc:
            with suppress(Exception):
                conn.rollback()
            raise DashboardCreatorError(
                "Failed to create dashboard",
                code="DASHBOARD_CREATE_FAILED",
                details={"error": str(exc)},
            ) from exc
        finally:
            self._close_connection(conn, cursor)

    def update_dashboard(
        self,
        dashboard_id: int,
        payload: Dict[str, Any],
        username: str = "system",
        owner_user_id: Optional[int] = None,
    ) -> Dict[str, Any]:
        existing = self.get_dashboard(dashboard_id)
        dash_name = self._normalize_dashboard_name(payload.get("dashboardName") or existing.get("dashboardName"))

        description = self._normalize_description(payload.get("description", existing.get("description")))
        payload_owner_user_id = self._to_int(payload.get("ownerUserId"))
        if payload_owner_user_id is not None:
            resolved_owner_user_id = payload_owner_user_id
        elif owner_user_id is not None:
            resolved_owner_user_id = owner_user_id
        else:
            resolved_owner_user_id = self._to_int(existing.get("ownerUserId"))
        is_active = self._to_flag(payload.get("isActive", existing.get("isActive", True)))

        conn, cursor, db_type, tables = self._open_connection()
        try:
            if db_type == "ORACLE":
                update_query = f"""
                    UPDATE {tables['dash_def']}
                    SET DASHNM = :dashnm,
                        DSCRPTN = :dscrptn,
                        OWNER_USRID = :owner_usrid,
                        IS_ACTV = :is_actv,
                        UPDTDBY = :updtdby,
                        UPDTDT = SYSTIMESTAMP
                    WHERE DASHID = :dashid
                      AND CURFLG = 'Y'
                """
                params = {
                    "dashnm": dash_name,
                    "dscrptn": description,
                    "owner_usrid": resolved_owner_user_id,
                    "is_actv": is_active,
                    "updtdby": username,
                    "dashid": dashboard_id,
                }
                self._execute(cursor, update_query, params)
            else:
                update_query = f"""
                    UPDATE {tables['dash_def']}
                    SET DASHNM = %s,
                        DSCRPTN = %s,
                        OWNER_USRID = %s,
                        IS_ACTV = %s,
                        UPDTDBY = %s,
                        UPDTDT = CURRENT_TIMESTAMP
                    WHERE DASHID = %s
                      AND CURFLG = 'Y'
                """
                params = (dash_name, description, resolved_owner_user_id, is_active, username, dashboard_id)
                self._execute(cursor, update_query, params)

            if "widgets" in payload:
                if db_type == "ORACLE":
                    deactivate_query = f"""
                        UPDATE {tables['dash_widget']}
                        SET CURFLG = 'N',
                            IS_ACTV = 'N',
                            UPDTDBY = :updtdby,
                            UPDTDT = SYSTIMESTAMP
                        WHERE DASHID = :dashid
                          AND CURFLG = 'Y'
                    """
                    self._execute(cursor, deactivate_query, {"updtdby": username, "dashid": dashboard_id})
                else:
                    deactivate_query = f"""
                        UPDATE {tables['dash_widget']}
                        SET CURFLG = 'N',
                            IS_ACTV = 'N',
                            UPDTDBY = %s,
                            UPDTDT = CURRENT_TIMESTAMP
                        WHERE DASHID = %s
                          AND CURFLG = 'Y'
                    """
                    self._execute(cursor, deactivate_query, (username, dashboard_id))

                widgets = self._normalize_widgets(payload.get("widgets") or [])
                for index, widget in enumerate(widgets, start=1):
                    self._insert_widget(cursor, db_type, tables, dashboard_id, widget, username, index)

            with suppress(Exception):
                conn.commit()
            return self.get_dashboard(dashboard_id)
        except DashboardCreatorError:
            with suppress(Exception):
                conn.rollback()
            raise
        except Exception as exc:
            with suppress(Exception):
                conn.rollback()
            raise DashboardCreatorError(
                "Failed to update dashboard",
                code="DASHBOARD_UPDATE_FAILED",
                details={"error": str(exc)},
            ) from exc
        finally:
            self._close_connection(conn, cursor)

    def delete_dashboard(self, dashboard_id: int, username: str = "system") -> Dict[str, Any]:
        self.get_dashboard(dashboard_id)
        conn, cursor, db_type, tables = self._open_connection()
        try:
            if db_type == "ORACLE":
                self._execute(
                    cursor,
                    f"""
                    UPDATE {tables['dash_widget']}
                    SET CURFLG = 'N', IS_ACTV = 'N', UPDTDBY = :updtdby, UPDTDT = SYSTIMESTAMP
                    WHERE DASHID = :dashid AND CURFLG = 'Y'
                    """,
                    {"updtdby": username, "dashid": dashboard_id},
                )
                self._execute(
                    cursor,
                    f"""
                    UPDATE {tables['dash_def']}
                    SET CURFLG = 'N', IS_ACTV = 'N', UPDTDBY = :updtdby, UPDTDT = SYSTIMESTAMP
                    WHERE DASHID = :dashid AND CURFLG = 'Y'
                    """,
                    {"updtdby": username, "dashid": dashboard_id},
                )
            else:
                self._execute(
                    cursor,
                    f"""
                    UPDATE {tables['dash_widget']}
                    SET CURFLG = 'N', IS_ACTV = 'N', UPDTDBY = %s, UPDTDT = CURRENT_TIMESTAMP
                    WHERE DASHID = %s AND CURFLG = 'Y'
                    """,
                    (username, dashboard_id),
                )
                self._execute(
                    cursor,
                    f"""
                    UPDATE {tables['dash_def']}
                    SET CURFLG = 'N', IS_ACTV = 'N', UPDTDBY = %s, UPDTDT = CURRENT_TIMESTAMP
                    WHERE DASHID = %s AND CURFLG = 'Y'
                    """,
                    (username, dashboard_id),
                )

            with suppress(Exception):
                conn.commit()
            return {"dashboardId": dashboard_id, "deleted": True}
        except Exception as exc:
            with suppress(Exception):
                conn.rollback()
            raise DashboardCreatorError(
                "Failed to delete dashboard",
                code="DASHBOARD_DELETE_FAILED",
                details={"error": str(exc)},
            ) from exc
        finally:
            self._close_connection(conn, cursor)

    def list_sql_sources(self) -> List[Dict[str, Any]]:
        try:
            return self.report_service.list_sql_sources()
        except ReportServiceError as exc:
            raise DashboardCreatorError(
                exc.message,
                status_code=exc.status_code,
                code=exc.code,
                details=exc.details,
            ) from exc

    def describe_sql_columns(self, sql_text: Optional[str], db_connection_id: Optional[int]) -> Dict[str, Any]:
        try:
            return self.report_service.describe_sql_columns(sql_text, db_connection_id)
        except ReportServiceError as exc:
            raise DashboardCreatorError(
                exc.message,
                status_code=exc.status_code,
                code=exc.code,
                details=exc.details,
            ) from exc

    def preview_widget_sql(
        self,
        sql_text: Optional[str],
        db_connection_id: Optional[int] = None,
        row_limit: Optional[int] = 100,
    ) -> Dict[str, Any]:
        if not sql_text or not str(sql_text).strip():
            raise DashboardCreatorError("sqlText is required", code="SQL_TEXT_REQUIRED")
        self._validate_read_only_sql(str(sql_text))

        try:
            limit_value = int(row_limit or 100)
        except Exception:
            limit_value = 100
        limit_value = max(1, min(limit_value, 1000))

        try:
            result = self.report_service._run_preview_query(
                connection_id=self._to_int(db_connection_id),
                sql_text=str(sql_text),
                row_limit=limit_value,
                parameters={},
            )
            rows = result.get("rows", [])
            safe_rows: List[Dict[str, Any]] = []
            for row in rows:
                safe_rows.append({key: self._json_safe(value) for key, value in row.items()})

            return {
                "rowLimit": limit_value,
                "rowCount": len(safe_rows),
                "columns": result.get("columns", []),
                "rows": safe_rows,
                "sourceDbType": result.get("dbType"),
            }
        except ReportServiceError as exc:
            raise DashboardCreatorError(
                exc.message,
                status_code=exc.status_code,
                code=exc.code,
                details=exc.details,
            ) from exc

    def _safe_file_name(self, value: str) -> str:
        cleaned = "".join(char if char.isalnum() or char in "-_" else "_" for char in (value or "dashboard"))
        return cleaned or "dashboard"

    def _resolve_widget_sql(self, widget: Dict[str, Any]) -> Dict[str, Any]:
        adhoc_sql = widget.get("adhocSql")
        if adhoc_sql and str(adhoc_sql).strip():
            return {
                "sqlText": str(adhoc_sql),
                "connectionId": self._to_int(widget.get("dbConnectionId")),
            }

        sql_source_id = self._to_int(widget.get("sqlSourceId"))
        if not sql_source_id:
            raise DashboardCreatorError(
                "Widget SQL source is missing",
                code="WIDGET_SQL_SOURCE_MISSING",
            )

        conn, cursor, db_type, tables = self.report_service._open_connection()
        try:
            source = self.report_service._load_sql_source_record(
                cursor=cursor,
                tables=tables,
                db_type=db_type,
                sql_source_id=sql_source_id,
                fallback_connection_id=self._to_int(widget.get("dbConnectionId")),
            )
            return {
                "sqlText": source.get("sqlText") or "",
                "connectionId": self._to_int(widget.get("dbConnectionId")) or source.get("connectionId"),
            }
        finally:
            self.report_service._close_connection(conn, cursor)

    def _collect_export_data(self, dashboard_id: int, row_limit: int = 500) -> Dict[str, Any]:
        dashboard = self.get_dashboard(dashboard_id)
        widgets = dashboard.get("widgets") or []
        if not widgets:
            raise DashboardCreatorError("Dashboard has no widgets to export", code="NO_WIDGETS_TO_EXPORT")

        export_widgets: List[Dict[str, Any]] = []
        for widget in widgets:
            sql_payload = self._resolve_widget_sql(widget)
            if not (sql_payload.get("sqlText") or "").strip():
                continue

            result = self.report_service._run_preview_query(
                connection_id=self._to_int(sql_payload.get("connectionId")),
                sql_text=str(sql_payload.get("sqlText")),
                row_limit=max(1, min(int(row_limit), 2000)),
                parameters={},
            )

            rows = result.get("rows", [])
            safe_rows = [
                {column: self._json_safe(value) for column, value in row.items()}
                for row in rows
            ]

            export_widgets.append(
                {
                    "widgetName": widget.get("widgetName") or "Widget",
                    "widgetType": widget.get("widgetType") or "TABLE",
                    "columns": result.get("columns", []),
                    "rows": safe_rows,
                    "rowCount": len(safe_rows),
                    "sourceDbType": result.get("dbType"),
                }
            )

        if not export_widgets:
            raise DashboardCreatorError("No widget datasets generated for export", code="NO_WIDGET_DATA")

        return {
            "dashboard": dashboard,
            "widgets": export_widgets,
        }

    def _render_pdf_bytes(self, export_payload: Dict[str, Any]) -> bytes:
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.graphics.charts.barcharts import VerticalBarChart
            from reportlab.graphics.charts.linecharts import HorizontalLineChart
            from reportlab.graphics.charts.piecharts import Pie
            from reportlab.graphics.shapes import Drawing, String
            from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle, PageBreak
        except ImportError as exc:
            raise DashboardCreatorError(
                "PDF export requires reportlab package. Install with: pip install reportlab",
                status_code=400,
                code="PDF_DEPENDENCY_MISSING",
            ) from exc

        dashboard = export_payload["dashboard"]
        widgets = export_payload["widgets"]

        output = io.BytesIO()
        document = SimpleDocTemplate(output, pagesize=landscape(letter))
        styles = getSampleStyleSheet()
        story = []

        title = dashboard.get("dashboardName") or "Dashboard Export"
        description = dashboard.get("description") or ""
        story.append(Paragraph(title, styles["Title"]))
        if description:
            story.append(Paragraph(description, styles["Normal"]))
        story.append(Spacer(1, 12))

        max_cols_per_page = 10
        max_rows_per_page = 25
        chart_palette = [
            colors.HexColor("#2563EB"),
            colors.HexColor("#0EA5E9"),
            colors.HexColor("#10B981"),
            colors.HexColor("#F59E0B"),
            colors.HexColor("#EF4444"),
            colors.HexColor("#8B5CF6"),
            colors.HexColor("#14B8A6"),
            colors.HexColor("#F97316"),
        ]

        def _to_float(value: Any) -> Optional[float]:
            try:
                if value is None or value == "":
                    return None
                return float(value)
            except (TypeError, ValueError):
                return None

        def _first_numeric_column(columns: List[str], rows: List[Dict[str, Any]]) -> Optional[str]:
            for column in columns:
                for row in rows:
                    number_value = _to_float(row.get(column))
                    if number_value is not None:
                        return column
            return None

        def _build_chart_drawing(widget_type: str, columns: List[str], rows: List[Dict[str, Any]]) -> Optional[Drawing]:
            if not columns or not rows:
                return None

            x_column = columns[0]
            y_column = _first_numeric_column(columns[1:] if len(columns) > 1 else columns, rows)
            if not y_column:
                return None

            sample_rows = rows[:20]
            labels = [str(row.get(x_column, ""))[:24] for row in sample_rows]
            values = [(_to_float(row.get(y_column)) or 0.0) for row in sample_rows]

            drawing = Drawing(720, 280)
            drawing.add(String(12, 262, f"{widget_type} chart: {y_column} by {x_column}", fontSize=10, fillColor=colors.HexColor("#334155")))

            if widget_type == "BAR":
                chart = VerticalBarChart()
                chart.x = 50
                chart.y = 45
                chart.height = 200
                chart.width = 620
                chart.data = [values]
                chart.categoryAxis.categoryNames = labels
                chart.valueAxis.valueMin = 0
                chart.barWidth = 10
                chart.bars[0].fillColor = chart_palette[0]
                chart.bars[0].strokeColor = colors.HexColor("#1E40AF")
                chart.categoryAxis.labels.fontSize = 7
                chart.categoryAxis.labels.angle = 25
                chart.categoryAxis.labels.dy = -8
                chart.valueAxis.labels.fontSize = 8
                chart.valueAxis.gridStrokeColor = colors.HexColor("#CBD5E1")
                drawing.add(chart)
                return drawing

            if widget_type in {"LINE", "AREA"}:
                chart = HorizontalLineChart()
                chart.x = 50
                chart.y = 45
                chart.height = 200
                chart.width = 620
                chart.data = [values]
                chart.categoryAxis.categoryNames = labels
                chart.valueAxis.valueMin = 0
                chart.lines[0].strokeColor = chart_palette[0]
                chart.lines[0].strokeWidth = 2
                chart.categoryAxis.labels.fontSize = 7
                chart.categoryAxis.labels.angle = 25
                chart.categoryAxis.labels.dy = -8
                chart.valueAxis.labels.fontSize = 8
                chart.valueAxis.gridStrokeColor = colors.HexColor("#CBD5E1")
                drawing.add(chart)
                return drawing

            if widget_type == "PIE":
                chart = Pie()
                chart.x = 220
                chart.y = 30
                chart.width = 280
                chart.height = 220
                chart.labels = labels
                chart.data = values
                chart.slices.strokeColor = colors.white
                chart.slices.strokeWidth = 0.5
                for slice_index in range(min(len(values), len(chart_palette))):
                    chart.slices[slice_index].fillColor = chart_palette[slice_index]
                drawing.add(chart)
                return drawing

            return None

        for index, widget in enumerate(widgets, start=1):
            columns = widget.get("columns", [])
            rows = widget.get("rows", [])
            row_count = int(widget.get("rowCount") or len(rows) or 0)
            source_db = widget.get("sourceDbType") or "-"
            widget_name = widget.get("widgetName") or f"Widget {index}"
            widget_type = str(widget.get("widgetType") or "TABLE").upper()

            story.append(Paragraph(f"{index}. {widget_name} ({widget_type})", styles["Heading3"]))
            story.append(Paragraph(f"Rows: {row_count} | Source DB: {source_db}", styles["Normal"]))

            if widget_type in {"BAR", "LINE", "AREA", "PIE"}:
                chart_drawing = _build_chart_drawing(widget_type, columns, rows)
                if chart_drawing:
                    story.append(chart_drawing)
                else:
                    story.append(Paragraph("Unable to render chart for this widget. Showing tabular fallback.", styles["Italic"]))

                if chart_drawing and index != len(widgets):
                    story.append(PageBreak())
                    continue

            if widget_type == "KPI":
                metric_column = _first_numeric_column(columns, rows)
                if metric_column:
                    total_value = sum(_to_float(row.get(metric_column)) or 0.0 for row in rows)
                    kpi_table = Table(
                        [[f"KPI ({metric_column})"], [f"{total_value:,.2f}"]],
                        colWidths=[740],
                    )
                    kpi_table.setStyle(
                        TableStyle(
                            [
                                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E2E8F0")),
                                ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#1E293B")),
                                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                                ("FONTSIZE", (0, 0), (-1, 0), 12),
                                ("BACKGROUND", (0, 1), (-1, 1), colors.HexColor("#F8FAFC")),
                                ("TEXTCOLOR", (0, 1), (-1, 1), colors.HexColor("#0F172A")),
                                ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold"),
                                ("FONTSIZE", (0, 1), (-1, 1), 24),
                                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
                                ("TOPPADDING", (0, 0), (-1, -1), 8),
                                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                            ]
                        )
                    )
                    story.append(kpi_table)
                else:
                    story.append(Paragraph("No numeric data available for KPI widget.", styles["Italic"]))

                if index != len(widgets):
                    story.append(PageBreak())
                continue

            if columns:
                column_chunks = [
                    columns[pos:pos + max_cols_per_page]
                    for pos in range(0, len(columns), max_cols_per_page)
                ]
                row_chunks = [
                    rows[pos:pos + max_rows_per_page]
                    for pos in range(0, len(rows), max_rows_per_page)
                ] or [[]]

                total_parts = max(1, len(column_chunks) * len(row_chunks))
                part_number = 1

                for col_chunk_index, col_chunk in enumerate(column_chunks):
                    for row_chunk_index, row_chunk in enumerate(row_chunks):
                        if total_parts > 1:
                            story.append(Paragraph(f"Part {part_number}/{total_parts}", styles["Italic"]))

                        row_start = (row_chunk_index * max_rows_per_page) + 1 if row_count else 0
                        row_end = (row_chunk_index * max_rows_per_page) + len(row_chunk)
                        col_start = (col_chunk_index * max_cols_per_page) + 1
                        col_end = (col_chunk_index * max_cols_per_page) + len(col_chunk)

                        if row_count:
                            summary = (
                                f"Rows: {row_count} | Showing rows {row_start}-{row_end} | "
                                f"Columns {col_start}-{col_end} | Source DB: {source_db}"
                            )
                        else:
                            summary = f"Rows: 0 | Columns {col_start}-{col_end} | Source DB: {source_db}"
                        story.append(Paragraph(summary, styles["Normal"]))

                        table_data = [col_chunk]
                        for row in row_chunk:
                            table_data.append([str(row.get(column, "")) for column in col_chunk])

                        table = Table(table_data, repeatRows=1)
                        table.setStyle(
                            TableStyle(
                                [
                                    ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                                    ("FONTSIZE", (0, 0), (-1, 0), 9),
                                    ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
                                    ("FONTSIZE", (0, 1), (-1, -1), 8),
                                ]
                            )
                        )
                        story.append(table)

                        is_last_part = part_number == total_parts
                        is_last_widget = index == len(widgets)
                        if not is_last_part or not is_last_widget:
                            story.append(PageBreak())

                        part_number += 1
            else:
                story.append(Paragraph("No columns available.", styles["Italic"]))
                if index != len(widgets):
                    story.append(PageBreak())

        document.build(story)
        output.seek(0)
        return output.getvalue()

    def _render_ppt_bytes(self, export_payload: Dict[str, Any]) -> bytes:
        try:
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise DashboardCreatorError(
                "PPT export requires python-pptx package. Install with: pip install python-pptx",
                status_code=400,
                code="PPT_DEPENDENCY_MISSING",
            ) from exc

        dashboard = export_payload["dashboard"]
        widgets = export_payload["widgets"]

        presentation = Presentation()

        title_slide_layout = presentation.slide_layouts[0]
        title_slide = presentation.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = dashboard.get("dashboardName") or "Dashboard Export"
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"Generated at {datetime.utcnow().isoformat()} UTC"

        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        def _to_float(value: Any) -> Optional[float]:
            try:
                if value is None or value == "":
                    return None
                return float(value)
            except (TypeError, ValueError):
                return None

        def _first_numeric_column(columns: List[str], rows: List[Dict[str, Any]]) -> Optional[str]:
            for column in columns:
                for row in rows:
                    number_value = _to_float(row.get(column))
                    if number_value is not None:
                        return column
            return None

        margin = Inches(0.4)
        content_left = margin
        content_width = slide_width - (2 * margin)
        summary_top = Inches(1.0)
        summary_height = Inches(0.45)
        table_top = summary_top + summary_height + Inches(0.1)
        table_height = slide_height - table_top - margin

        max_cols_per_slide = 6
        max_rows_per_slide = 18

        for widget in widgets:
            widget_name = widget.get("widgetName") or "Widget"
            widget_type = str(widget.get("widgetType") or "TABLE").upper()
            columns = widget.get("columns", [])
            rows = widget.get("rows", [])
            source_db = widget.get("sourceDbType") or "-"
            total_row_count = int(widget.get("rowCount") or len(rows) or 0)

            chart_type_map = {
                "BAR": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "LINE": XL_CHART_TYPE.LINE_MARKERS,
                "AREA": XL_CHART_TYPE.AREA,
                "PIE": XL_CHART_TYPE.PIE,
            }

            x_column = columns[0] if columns else None
            y_column = _first_numeric_column(columns[1:] if len(columns) > 1 else columns, rows)

            if widget_type in chart_type_map and x_column and y_column and rows:
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                slide.shapes.title.text = f"{widget_name} ({widget_type})"

                textbox = slide.shapes.add_textbox(content_left, summary_top, content_width, summary_height)
                text_frame = textbox.text_frame
                text_frame.text = f"Rows: {total_row_count} | Source DB: {source_db}"
                text_frame.paragraphs[0].font.size = Pt(11)

                chart_labels = []
                chart_values = []
                for row in rows[:30]:
                    chart_labels.append(str(row.get(x_column, "")))
                    chart_values.append(_to_float(row.get(y_column)) or 0.0)

                chart_data = CategoryChartData()
                chart_data.categories = chart_labels
                chart_data.add_series(y_column, chart_values)

                chart_shape = slide.shapes.add_chart(
                    chart_type_map[widget_type],
                    content_left,
                    table_top,
                    content_width,
                    table_height,
                    chart_data,
                )
                chart = chart_shape.chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
            elif widget_type == "KPI" and rows and columns:
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                slide.shapes.title.text = f"{widget_name} ({widget_type})"

                textbox = slide.shapes.add_textbox(content_left, summary_top, content_width, summary_height)
                text_frame = textbox.text_frame
                text_frame.text = f"Rows: {total_row_count} | Source DB: {source_db}"
                text_frame.paragraphs[0].font.size = Pt(11)

                metric_column = _first_numeric_column(columns, rows)
                if metric_column:
                    total_value = sum(_to_float(row.get(metric_column)) or 0.0 for row in rows)
                    kpi_box = slide.shapes.add_textbox(content_left, table_top, content_width, table_height)
                    kpi_frame = kpi_box.text_frame
                    kpi_frame.clear()
                    metric_title = kpi_frame.paragraphs[0]
                    metric_title.text = f"KPI ({metric_column})"
                    metric_title.font.size = Pt(18)

                    metric_value = kpi_frame.add_paragraph()
                    metric_value.text = f"{total_value:,.2f}"
                    metric_value.font.size = Pt(42)
                else:
                    no_data_box = slide.shapes.add_textbox(content_left, table_top, content_width, Inches(1.0))
                    no_data_frame = no_data_box.text_frame
                    no_data_frame.text = "No numeric data available for KPI widget."
                    no_data_frame.paragraphs[0].font.size = Pt(12)
            elif columns:
                column_chunks = [
                    columns[index:index + max_cols_per_slide]
                    for index in range(0, len(columns), max_cols_per_slide)
                ]
                row_chunks = [
                    rows[index:index + max_rows_per_slide]
                    for index in range(0, len(rows), max_rows_per_slide)
                ] or [[]]

                total_parts = max(1, len(column_chunks) * len(row_chunks))
                part_number = 1

                for col_chunk_index, col_chunk in enumerate(column_chunks):
                    for row_chunk_index, row_chunk in enumerate(row_chunks):
                        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                        title_text = f"{widget_name} ({widget_type})"
                        if total_parts > 1:
                            title_text = f"{title_text} - Part {part_number}/{total_parts}"
                        slide.shapes.title.text = title_text

                        row_start = (row_chunk_index * max_rows_per_slide) + 1 if rows else 0
                        row_end = (row_chunk_index * max_rows_per_slide) + len(row_chunk)
                        col_start = (col_chunk_index * max_cols_per_slide) + 1
                        col_end = (col_chunk_index * max_cols_per_slide) + len(col_chunk)

                        textbox = slide.shapes.add_textbox(content_left, summary_top, content_width, summary_height)
                        text_frame = textbox.text_frame
                        if rows:
                            text_frame.text = (
                                f"Rows: {total_row_count} | Showing rows {row_start}-{row_end} | "
                                f"Columns {col_start}-{col_end} | Source DB: {source_db}"
                            )
                        else:
                            text_frame.text = f"Rows: 0 | Columns {col_start}-{col_end} | Source DB: {source_db}"
                        text_frame.paragraphs[0].font.size = Pt(10)

                        table_shape = slide.shapes.add_table(
                            rows=len(row_chunk) + 1,
                            cols=len(col_chunk),
                            left=content_left,
                            top=table_top,
                            width=content_width,
                            height=table_height,
                        )
                        table = table_shape.table

                        for column in table.columns:
                            column.width = int(content_width / max(len(col_chunk), 1))

                        for column_index, column_name in enumerate(col_chunk):
                            header_cell = table.cell(0, column_index)
                            header_cell.text = str(column_name)
                            if header_cell.text_frame.paragraphs:
                                header_cell.text_frame.paragraphs[0].font.size = Pt(9)

                        for row_index, row in enumerate(row_chunk, start=1):
                            for column_index, column_name in enumerate(col_chunk):
                                value = str(row.get(column_name, ""))
                                max_len = 80
                                table.cell(row_index, column_index).text = value[:max_len]
                                if table.cell(row_index, column_index).text_frame.paragraphs:
                                    table.cell(row_index, column_index).text_frame.paragraphs[0].font.size = Pt(8)

                        part_number += 1
            else:
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                slide.shapes.title.text = f"{widget_name} ({widget_type})"

                textbox = slide.shapes.add_textbox(content_left, summary_top, content_width, summary_height)
                text_frame = textbox.text_frame
                text_frame.text = f"Rows: {total_row_count} | Source DB: {source_db}"
                text_frame.paragraphs[0].font.size = Pt(11)

                no_data_box = slide.shapes.add_textbox(content_left, table_top, content_width, Inches(1.0))
                no_data_frame = no_data_box.text_frame
                no_data_frame.text = "No tabular data available for this widget."
                no_data_frame.paragraphs[0].font.size = Pt(12)

        output = io.BytesIO()
        presentation.save(output)
        output.seek(0)
        return output.getvalue()

    def _log_export(
        self,
        dashboard_id: int,
        export_format: str,
        username: str,
        status: str,
        message: Optional[str],
        file_name: Optional[str],
        file_size_bytes: Optional[int],
    ):
        conn, cursor, db_type, tables = self._open_connection()
        try:
            if db_type == "ORACLE":
                export_id = self._next_id(cursor, "DMS_DASH_EXPORT_LOG_SEQ")
                query = f"""
                    INSERT INTO {tables['dash_export_log']} (
                        EXPID, DASHID, EXPRT_FMT, EXPRT_BY, EXPRT_AT,
                        STTS, MSG, FILE_NM, FILE_SZ_BYTES
                    ) VALUES (
                        :expid, :dashid, :exprt_fmt, :exprt_by, SYSTIMESTAMP,
                        :stts, :msg, :file_nm, :file_sz_bytes
                    )
                """
                params = {
                    "expid": export_id,
                    "dashid": dashboard_id,
                    "exprt_fmt": export_format,
                    "exprt_by": username,
                    "stts": status,
                    "msg": message,
                    "file_nm": file_name,
                    "file_sz_bytes": file_size_bytes,
                }
                self._execute(cursor, query, params)
            else:
                query = f"""
                    INSERT INTO {tables['dash_export_log']} (
                        DASHID, EXPRT_FMT, EXPRT_BY, EXPRT_AT,
                        STTS, MSG, FILE_NM, FILE_SZ_BYTES
                    ) VALUES (
                        %s, %s, %s, CURRENT_TIMESTAMP,
                        %s, %s, %s, %s
                    )
                """
                params = (
                    dashboard_id,
                    export_format,
                    username,
                    status,
                    message,
                    file_name,
                    file_size_bytes,
                )
                self._execute(cursor, query, params)

            with suppress(Exception):
                conn.commit()
        except Exception:
            with suppress(Exception):
                conn.rollback()
        finally:
            self._close_connection(conn, cursor)

    def _serialize_export_history_row(self, row: Dict[str, Any]) -> Dict[str, Any]:
        return {
            "exportId": self._to_int(row.get("expid")),
            "dashboardId": self._to_int(row.get("dashid")),
            "dashboardName": row.get("dashnm"),
            "exportFormat": row.get("exprt_fmt"),
            "exportedBy": row.get("exprt_by"),
            "exportedAt": self._to_iso(row.get("exprt_at")),
            "status": row.get("stts"),
            "message": self._read_lob(row.get("msg")),
            "fileName": row.get("file_nm"),
            "fileSizeBytes": self._to_int(row.get("file_sz_bytes")),
        }

    def list_export_history(
        self,
        dashboard_id: Optional[int] = None,
        limit: int = 50,
    ) -> List[Dict[str, Any]]:
        try:
            effective_limit = int(limit)
        except Exception:
            effective_limit = 50
        effective_limit = max(1, min(effective_limit, 500))

        conn, cursor, db_type, tables = self._open_connection()
        try:
            builder = _ParamBuilder(db_type)
            query = f"""
                SELECT
                    e.EXPID AS expid,
                    e.DASHID AS dashid,
                    d.DASHNM AS dashnm,
                    e.EXPRT_FMT AS exprt_fmt,
                    e.EXPRT_BY AS exprt_by,
                    e.EXPRT_AT AS exprt_at,
                    e.STTS AS stts,
                    e.MSG AS msg,
                    e.FILE_NM AS file_nm,
                    e.FILE_SZ_BYTES AS file_sz_bytes
                FROM {tables['dash_export_log']} e
                JOIN {tables['dash_def']} d ON d.DASHID = e.DASHID
                WHERE d.CURFLG = 'Y'
            """

            if dashboard_id is not None:
                query += f" AND e.DASHID = {builder.add(dashboard_id, 'dash')}"

            query += " ORDER BY e.EXPRT_AT DESC"

            if db_type == "ORACLE":
                query = f"SELECT * FROM ({query}) WHERE ROWNUM <= {builder.add(effective_limit, 'limit')}"
                self._execute(cursor, query, builder.params)
            else:
                query += f" LIMIT {builder.add(effective_limit, 'limit')}"
                self._execute(cursor, query, builder.params)

            rows = self._fetch_all_dict(cursor)
            return [self._serialize_export_history_row(row) for row in rows]
        finally:
            self._close_connection(conn, cursor)

    def export_dashboard(
        self,
        dashboard_id: int,
        export_format: str,
        username: str = "system",
        row_limit: int = 500,
    ) -> Dict[str, Any]:
        normalized_format = (export_format or "").upper()
        if normalized_format not in {"PDF", "PPT"}:
            raise DashboardCreatorError(
                f"Unsupported export format: {export_format}",
                code="UNSUPPORTED_EXPORT_FORMAT",
            )

        try:
            row_limit_value = int(row_limit)
        except Exception:
            row_limit_value = 500
        row_limit_value = max(1, min(row_limit_value, 2000))

        dashboard_name = "dashboard"
        try:
            export_payload = self._collect_export_data(dashboard_id, row_limit=row_limit_value)
            dashboard_name = export_payload["dashboard"].get("dashboardName") or "dashboard"

            if normalized_format == "PDF":
                content = self._render_pdf_bytes(export_payload)
                media_type = "application/pdf"
                extension = "pdf"
            else:
                content = self._render_ppt_bytes(export_payload)
                media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                extension = "pptx"

            safe_name = self._safe_file_name(dashboard_name)
            file_name = f"{safe_name}.{extension}"

            self._log_export(
                dashboard_id=dashboard_id,
                export_format=normalized_format,
                username=username,
                status="SUCCESS",
                message=None,
                file_name=file_name,
                file_size_bytes=len(content),
            )

            return {
                "fileName": file_name,
                "mediaType": media_type,
                "content": content,
                "size": len(content),
            }
        except DashboardCreatorError as exc:
            self._log_export(
                dashboard_id=dashboard_id,
                export_format=normalized_format,
                username=username,
                status="FAILED",
                message=exc.message,
                file_name=None,
                file_size_bytes=None,
            )
            raise
        except Exception as exc:
            self._log_export(
                dashboard_id=dashboard_id,
                export_format=normalized_format,
                username=username,
                status="FAILED",
                message=str(exc),
                file_name=None,
                file_size_bytes=None,
            )
            raise DashboardCreatorError(
                "Failed to export dashboard",
                code="DASHBOARD_EXPORT_FAILED",
                details={"error": str(exc)},
            ) from exc
